import json
from pptx import Presentation
from PIL import Image
import io
import pytesseract

class PPTXLoader:
    """
    Load a PowerPoint presentation and extract content from it.
    """
    def __init__(self, pptx_file):
        """
        Initialize the SlidesBuilder with a PowerPoint file.

        :param pptx_file: Path to the PowerPoint file
        """
        # Should auto apply the presentation layout to the new slides added
        self.presentation = Presentation(pptx_file)
    
    @staticmethod
    def extract_text_from_notes_between_delimiters(slide, start_delimiter="###", end_delimiter="###"):
        """
        Extract text segments from the notes of a specific slide that are between specified start and end delimiters.

        :param slide: The slide object from which to extract notes.
        :param start_delimiter: The start delimiter string to search for in the notes.
        :param end_delimiter: The end delimiter string that marks the end of the text segment.
        :return: A list of extracted text segments found between each pair of start and end delimiters.
        """
        extracted_texts = []

        # Check if the slide has notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text

            # Initialize search position
            search_pos = 0

            # Search for the first occurrence of the start delimiter
            start_pos = notes_text.find(start_delimiter, search_pos)

            while start_pos != -1:
                # Find the end delimiter starting from the character after the start delimiter
                end_pos = notes_text.find(end_delimiter, start_pos + len(start_delimiter))

                # If the end delimiter is found, extract the text between delimiters
                if end_pos != -1:
                    # Extract text segment and add to the list
                    extracted_text = notes_text[start_pos + len(start_delimiter):end_pos].strip()
                    extracted_texts.append(extracted_text)

                    # Update search position to continue searching after the current end delimiter
                    search_pos = end_pos + len(end_delimiter)
                else:
                    # If no end delimiter is found, stop searching
                    break

                # Search for the next occurrence of the start delimiter
                start_pos = notes_text.find(start_delimiter, search_pos)

        return extracted_texts
        
    def extract_text_from_slide_notes_after_delimiter(self, slide, delimiter="###"):
        """
        Extract text from the notes of a specific slide after a specific delimiter.

        :param slide: The slide object from which to extract notes.
        :param delimiter: The delimiter string to search for in the notes.
        :return: Extracted text after the delimiter, or an empty string if the delimiter is not found.
        """
        # Initialize the return text as an empty string
        extracted_text = ""

        # Check if the slide has notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text

            # Find the delimiter in the notes text
            delimiter_position = notes_text.find(delimiter)
            if delimiter_position != -1:
                # Extract text after the delimiter
                extracted_text = notes_text[delimiter_position + len(delimiter):].strip()

        return extracted_text

    def extract_notes_after_delimiter_for_all_slides(self, delimiter):
        """
        Extract text from the presentation notes after a specific delimiter for all slides.

        :param delimiter: The delimiter string to search for in the notes.
        :return: A dictionary with slide numbers as keys and extracted text as values.
        """
        extracted_text = {}

        for i, slide in enumerate(self.presentation.slides):
            # Use the extraction function for the specific slide
            text_after_delimiter = self.extract_text_from_slide_notes_after_delimiter(slide, delimiter)
            if text_after_delimiter:
                # Store the extracted text with the slide number as the key
                extracted_text[i + 1] = text_after_delimiter

        return extracted_text

    def extract_text_from_slide(slide):
        """
        Extract the text from a slide.
        
        :param slide: Slide to extract text from
        :return: Text from the slide
        """
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        return text.strip()

    def interpret_image(image):
        """
        Interpret an image and return a description of it.
        
        :param image: Image to interpret
        :return: Description of the image
        """
        
        # Convert the image to a format suitable for text recognition
        image = image.convert('L')  # Convert to grayscale

        # Use Tesseract to do OCR on the image
        text = pytesseract.image_to_string(image)

        return text

    def get_pptx_content(self):
        """
        Prosess a PowerPoint presentation and return a JSON string with the content of each slide.
        
        :return: JSON string with the content of each slide
        """
        presentation = self.presentation
        slides_content = {}

        for i, slide in enumerate(presentation.slides):
            slide_content = {"text": self.extract_text_from_slide(slide), "images": []}

            for shape in slide.shapes:
                if shape.shape_type == 13:  # This is the type for Picture
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    image_description = self.interpret_image(image)
                    slide_content["images"].append(image_description)

            slides_content[f"Slide {i + 1}"] = slide_content

        return json.dumps(slides_content, indent=4)
