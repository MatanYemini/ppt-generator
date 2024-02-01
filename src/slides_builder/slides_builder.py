from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from content_generators.qr_generator import QRGenerator
import json

class SlidesBuilder:
    def __init__(self, pptx_file):
        """
        Initialize the SlidesBuilder with a PowerPoint file.

        :param pptx_file: Path to the PowerPoint file
        """
        # Should auto apply the presentation layout to the new slides added
        self.presentation = Presentation(pptx_file)
        
    def inspect_slide_layouts(self):
        """
        Inspect the slide layouts in the presentation.
        
        :return: Information about the slide layouts in the presentation
        """
        layouts_info = {}

        for i, layout in enumerate(self.presentation.slide_layouts):
            layout_info = {'name': layout.name, 'placeholders': []}
            for placeholder in layout.placeholders:
                layout_info['placeholders'].append({
                    'index': placeholder.placeholder_format.idx,
                    'type': placeholder.placeholder_format.type
                })
            layouts_info[i] = layout_info

        return layouts_info

    def add_slide(self, layout_idx=1, title="New Slide", content=""):
        """
        Add a new slide to the presentation with a title and content.

        :param title: Title of the new slide
        :param content: Content of the new slide
        :return: The new slide
        """
        # Use the first slide layout by default (typically a title slide)
        slide_layout = self.presentation.slide_layouts[layout_idx]

        # Add a slide
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set the title and content if placeholders are available
        if slide.shapes.title:
            slide.shapes.title.text = title

        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = content
                
        return slide

    def add_slide_with_title_and_content(self, title, content, slide_index=None):
        """
        Add a slide with a title at the top and content following.

        :param title: Title of the slide
        :param content: Content of the slide
        :return: The new slide
        """
        slide_layout = self.presentation.slide_layouts[1]  # A layout with title and content placeholders
        
        # Add a slide at the specified index, if provided
        if slide_index is not None and 0 <= slide_index <= len(self.presentation.slides):
            # This is a placeholder for the desired functionality
            # Actual implementation would depend on extending python-pptx or using a workaround
            slide = self.insert_slide_at_index(slide_layout, slide_index)
        else:
            # Add the slide at the end if no index is specified or if the index is out of range
            slide = self.presentation.slides.add_slide(slide_layout)
        
        # Assign the title and content if the placeholders exist
        title_placeholder = slide.shapes.title
        content_placeholder = None

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == "BODY":
                content_placeholder = placeholder
                break

        if title_placeholder is not None:
            title_placeholder.text = title
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)

        if content_placeholder is not None:
            content_placeholder.text = content
        else:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            content_box.text_frame.text = content
            
        return slide

    def insert_slide_at_index(self, slide_layout, slide_index):
        """
        Insert a slide at a specific index. This is a workaround method that creates a new presentation,
        copies slides to it, and inserts the new slide at the desired position.

        :param slide_layout: The layout to use for the new slide.
        :param slide_index: The index at which to insert the new slide.
        :return: The newly inserted slide object.
        """
        
        # Ensure the slide index is within bounds
        slide_index = max(0, min(slide_index, len(self.presentation.slides)))

        # Insert the new slide at the slide_index
        new_slide = self.presentation.slides.add_slide(slide_layout)

        # The new slide that was created at the last index
        # Move the new slide to the desired index
        self.move_slide(len(self.presentation.slides) - 1, slide_index)

        return new_slide

    def xml_slides(self):
        return self.presentation.slides._sldIdLst  # pylint: disable=protected-access

    def move_slide(self, old_index, new_index):
        xml_slides_info = self.xml_slides()
        slides = list(xml_slides_info)
        xml_slides_info.remove(slides[old_index])
        xml_slides_info.insert(new_index, slides[old_index])


    def safe_add_slide_with_title_and_content(self, title, content, slide_index=None):
        """
        Add a slide with a title at the top and content following, using default positions if specific placeholders are not found.

        :param title: Title of the slide
        :param content: Content of the slide
        :param slide_index: Index of the slide to add
        :return: The new slide
        """
        try:
            # Attempt to use a layout with title and content placeholders
            slide_layout = self.presentation.slide_layouts[1]
        except IndexError:
            # If the specified layout does not exist, use the default blank layout
            slide_layout = self.presentation.slide_layouts[5]  # Typically, index 5 is a blank layout

        # Add a slide at the specified index, if provided
        if slide_index is not None and 0 <= slide_index <= len(self.presentation.slides):
            # This is a placeholder for the desired functionality
            # Actual implementation would depend on extending python-pptx or using a workaround
            slide = self.insert_slide_at_index(slide_layout, slide_index)
        else:
            # Add the slide at the end if no index is specified or if the index is out of range
            slide = self.presentation.slides.add_slide(slide_layout)
        
        # Attempt to assign the title and content if the placeholders exist
        title_placeholder_found = False
        content_placeholder_found = False

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == "TITLE":
                placeholder.text = title
                title_placeholder_found = True
            elif placeholder.placeholder_format.type == "BODY":
                placeholder.text = content
                content_placeholder_found = True

        # If title placeholder is not found, add a textbox for the title
        if not title_placeholder_found and title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)

        # If content placeholder is not found, add a textbox for the content
        if not content_placeholder_found and content:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.text = content

        return slide
                
    def add_qr_slide(self, qr_data, text_dict, slide_title, slide_content="", slide_index=None):
        """
        Add a slide with specified number of QR codes containing the given data
        and corresponding text next to each QR code.

        :param qr_data: Data to be encoded in the QR codes
        :param text_dict: Dictionary with index and text for each QR code
        :param slide_title: Title of the new slide
        :param slide_content: Content of the new slide
        :param slide_index: Index of the slide to add
        """
        slide = self.add_slide_with_title_and_content(slide_title, slide_content, slide_index)

        # Define size and spacing
        qr_size = Inches(0.6)  # Size of each QR code
        qr_margin = Inches(0.2)  # Margin between QR codes
        max_qr = min(len(text_dict), 4)  # Limit the number of QR codes to 4 or the number of items in text_dict

        # Assuming QR codes are placed in the content area
        content_top = Inches(1.5)  # Starting position for QR codes
        content_left = Inches(0.5)  # Left margin for QR codes

        for i in range(max_qr):
            # Calculate position for QR code
            top = content_top + i * (qr_size + qr_margin)

            # Generate and add QR code to slide
            qr_img = QRGenerator.generate_qr_code(json.dumps(qr_data))
            qr_bytes = BytesIO()
            qr_img.save(qr_bytes)
            qr_bytes.seek(0)

            slide.shapes.add_picture(qr_bytes, content_left, top, qr_size, qr_size)

            # Add a textbox for each QR code
            if i in text_dict:
                textbox = slide.shapes.add_textbox(content_left + qr_size + qr_margin, top,
                                                   Inches(3), qr_size)
                textbox.text = text_dict[i]
                
    def add_full_slide_image(self, image_path, layout_idx=6, slide_index=None):
        """
        Add a new slide to the presentation with an image covering the entire slide.

        :param image_path: Path to the image file
        :param layout_idx: Index of the slide layout to use, default is 6 for a blank slide
        :param slide_index: Index of the slide to add
        :return: The new slide
        """
        # Use a blank slide layout by default
        slide_layout = self.presentation.slide_layouts[layout_idx]

        # Add a slide at the specified index, if provided
        if slide_index is not None and 0 <= slide_index <= len(self.presentation.slides):
            # This is a placeholder for the desired functionality
            # Actual implementation would depend on extending python-pptx or using a workaround
            slide = self.insert_slide_at_index(slide_layout, slide_index)
        else:
            # Add the slide at the end if no index is specified or if the index is out of range
            slide = self.presentation.slides.add_slide(slide_layout)

        # Get slide dimensions
        slide_width = self.presentation.slide_width
        slide_height = self.presentation.slide_height

        # Add the image to cover the entire slide
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

        return slide
    
    def append_notes_to_slide(self, slide, notes, delimiter="###"):
        """
        Append notes to a slide's notes page, using a new paragraph to preserve existing formatting.

        :param slide: The slide object to which notes will be added
        :param notes: The notes text to be appended
        :param delimiter: A string delimiter to separate different sections of notes
        """

        # Access the text frame of the notes slide
        text_frame = slide.notes_slide.notes_text_frame

        
        # Add a new paragraph with delimiter and content
        new_paragraph = text_frame.add_paragraph()
        new_paragraph.text = delimiter
        new_paragraph.text += notes        
        new_paragraph.text += delimiter

    def save(self, output_file):
        """
        Save the updated presentation to a new file.

        :param output_file: Path to save the updated PowerPoint file
        """
        self.presentation.save(output_file)
    